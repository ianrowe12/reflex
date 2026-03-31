#!/usr/bin/env python3
"""Build Reflex Team Walkthrough PowerPoint from structured markdown slide files."""

import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Design Constants ──────────────────────────────────────────────────────────
BG_DARK_NAVY = RGBColor(0x1B, 0x2A, 0x4A)
BG_DIVIDER = RGBColor(0x24, 0x3B, 0x5E)
CLR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLR_GRAY = RGBColor(0xB0, 0xBE, 0xC5)
CLR_GOLD = RGBColor(0xD4, 0xA8, 0x43)
CLR_RED = RGBColor(0xC0, 0x39, 0x2B)
CLR_GREEN = RGBColor(0x27, 0xAE, 0x60)
CLR_BLUE = RGBColor(0x34, 0x98, 0xDB)
CLR_DEEP_REF = RGBColor(0x78, 0x90, 0x9C)
CLR_PURPLE = RGBColor(0x9B, 0x59, 0xB6)
CLR_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
CLR_ROW_ALT = RGBColor(0x22, 0x36, 0x56)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.75)
CONTENT_W = SLIDE_W - 2 * MARGIN
TITLE_TOP = Inches(0.4)
TITLE_H = Inches(0.7)
SUBTITLE_TOP = Inches(1.1)
SUBTITLE_H = Inches(0.5)
BODY_TOP = Inches(1.7)


# ── Slide Parsing ─────────────────────────────────────────────────────────────
def parse_slides(filepath):
    """Parse a markdown file into a list of slide dicts."""
    with open(filepath, "r") as f:
        text = f.read()
    raw_slides = re.split(r"---SLIDE---", text)[1:]  # skip preamble
    slides = []
    for raw in raw_slides:
        raw = raw.split("---END---")[0]
        slide = {}
        # type
        m = re.search(r'^type:\s*(.+)', raw, re.M)
        slide['type'] = m.group(1).strip() if m else 'content'
        # title
        m = re.search(r'^title:\s*"?(.+?)"?\s*$', raw, re.M)
        slide['title'] = m.group(1).strip().strip('"') if m else ''
        # subtitle
        m = re.search(r'^subtitle:\s*"?(.+?)"?\s*$', raw, re.M)
        slide['subtitle'] = m.group(1).strip().strip('"') if m else ''
        # bullets
        bullet_section = re.search(r'^bullets:\s*\n((?:\s*-\s*.+\n?)+)', raw, re.M)
        if bullet_section:
            slide['bullets'] = [
                re.sub(r'^\s*-\s*"?', '', b).rstrip().rstrip('"')
                for b in bullet_section.group(1).strip().split('\n')
                if b.strip() and b.strip() != '-'
            ]
        else:
            slide['bullets'] = []
        # speaker_notes
        m = re.search(r'^speaker_notes:\s*"(.+?)"(?:\s*$)', raw, re.M | re.S)
        slide['speaker_notes'] = m.group(1).strip() if m else ''
        # deep_dive
        m = re.search(r'^deep_dive:\s*"?(.+?)"?\s*$', raw, re.M)
        slide['deep_dive'] = m.group(1).strip().strip('"') if m else ''
        # left_column / right_column
        for col in ('left_column', 'right_column'):
            col_section = re.search(rf'^{col}:\s*\n((?:\s*-\s*.+\n?)+)', raw, re.M)
            if col_section:
                slide[col] = [
                    re.sub(r'^\s*-\s*"?', '', b).rstrip().rstrip('"')
                    for b in col_section.group(1).strip().split('\n')
                    if b.strip()
                ]
            else:
                slide[col] = []
        # table
        table_match = re.search(
            r'^table:\s*\n\s*headers:\s*\[(.+?)\]\s*\n\s*rows:\s*\n((?:\s*-\s*\[.+?\]\s*\n?)+)',
            raw, re.M
        )
        if table_match:
            headers = [h.strip().strip('"').strip("'") for h in table_match.group(1).split(',')]
            row_text = table_match.group(2)
            rows = []
            for row_match in re.finditer(r'\[(.+?)\]', row_text):
                cells = [c.strip().strip('"').strip("'") for c in row_match.group(1).split('",')]
                # Clean up remaining quotes
                cells = [c.strip().strip('"').strip("'") for c in cells]
                rows.append(cells)
            slide['table'] = {'headers': headers, 'rows': rows}
        else:
            slide['table'] = None
        # visual_description
        m = re.search(r'^visual_description:\s*"(.+?)"(?:\s*$)', raw, re.M | re.S)
        slide['visual_description'] = m.group(1).strip() if m else ''

        slides.append(slide)
    return slides


# ── Slide Building Helpers ────────────────────────────────────────────────────
def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=CLR_WHITE, align=PP_ALIGN.LEFT,
                italic=False, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.italic = italic
    p.alignment = align
    return txBox


def add_bullets(slide, left, top, width, height, items, font_size=18,
                color=CLR_WHITE, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Strip markdown bold markers for display
        clean = re.sub(r'\*\*(.+?)\*\*', r'\1', item)
        # Check if item starts with indent
        indent_level = 0
        if clean.startswith('  ') or clean.startswith('\t'):
            indent_level = 1
            clean = clean.strip()

        p.text = clean
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
        p.level = indent_level

        # Bold portions
        if '**' in item:
            p.clear()
            parts = re.split(r'(\*\*.+?\*\*)', item)
            for part in parts:
                run = p.add_run()
                if part.startswith('**') and part.endswith('**'):
                    run.text = part[2:-2]
                    run.font.bold = True
                else:
                    run.text = part
                run.font.size = Pt(font_size)
                run.font.color.rgb = color
    return txBox


def add_slide_number(slide, num):
    add_textbox(
        slide,
        SLIDE_W - Inches(0.75), SLIDE_H - Inches(0.4),
        Inches(0.5), Inches(0.3),
        str(num), font_size=10, color=CLR_GRAY, align=PP_ALIGN.RIGHT
    )


def add_deep_dive(slide, ref):
    if ref:
        add_textbox(
            slide,
            MARGIN, SLIDE_H - Inches(0.5),
            Inches(4), Inches(0.3),
            f"Deep dive: {ref}",
            font_size=12, color=CLR_DEEP_REF, italic=True
        )


def add_speaker_notes(slide, notes):
    if notes:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes


def add_title_bar(slide, title, subtitle=''):
    add_textbox(
        slide, MARGIN, TITLE_TOP, CONTENT_W, TITLE_H,
        title, font_size=28, bold=True, color=CLR_WHITE
    )
    if subtitle:
        add_textbox(
            slide, MARGIN, SUBTITLE_TOP, CONTENT_W, SUBTITLE_H,
            subtitle, font_size=20, color=CLR_GRAY
        )


# ── Slide Type Builders ───────────────────────────────────────────────────────
def build_title_slide(prs, data, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, BG_DARK_NAVY)
    # Large centered title
    add_textbox(
        slide, Inches(1), Inches(2), Inches(11.333), Inches(1.2),
        data['title'], font_size=36, bold=True, color=CLR_WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
    )
    # Subtitle
    if data['subtitle']:
        add_textbox(
            slide, Inches(1), Inches(3.2), Inches(11.333), Inches(0.6),
            data['subtitle'], font_size=20, color=CLR_GRAY,
            align=PP_ALIGN.CENTER
        )
    # Tagline from bullets
    if data['bullets']:
        add_textbox(
            slide, Inches(2), Inches(4.0), Inches(9.333), Inches(0.5),
            data['bullets'][0], font_size=16, color=CLR_GOLD,
            align=PP_ALIGN.CENTER, italic=True
        )
    # Date at bottom
    add_textbox(
        slide, Inches(1), Inches(6.2), Inches(11.333), Inches(0.4),
        "March 2026", font_size=14, color=CLR_GRAY, align=PP_ALIGN.CENTER
    )
    add_slide_number(slide, num)
    add_speaker_notes(slide, data['speaker_notes'])
    add_deep_dive(slide, data.get('deep_dive', ''))


def build_divider_slide(prs, data, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DIVIDER)
    add_textbox(
        slide, Inches(1), Inches(2.5), Inches(11.333), Inches(1),
        data['title'], font_size=32, bold=True, color=CLR_GOLD,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
    )
    if data['subtitle']:
        add_textbox(
            slide, Inches(1), Inches(3.7), Inches(11.333), Inches(0.6),
            data['subtitle'], font_size=20, color=CLR_GRAY,
            align=PP_ALIGN.CENTER
        )
    add_slide_number(slide, num)
    add_speaker_notes(slide, data['speaker_notes'])
    add_deep_dive(slide, data.get('deep_dive', ''))


def build_content_slide(prs, data, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK_NAVY)
    add_title_bar(slide, data['title'], data.get('subtitle', ''))
    body_top = BODY_TOP if not data.get('subtitle') else Inches(1.9)
    body_h = SLIDE_H - body_top - Inches(0.8)
    if data['bullets']:
        add_bullets(slide, MARGIN, body_top, CONTENT_W, body_h, data['bullets'][:6])
    add_slide_number(slide, num)
    add_speaker_notes(slide, data['speaker_notes'])
    add_deep_dive(slide, data.get('deep_dive', ''))


def build_table_slide(prs, data, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK_NAVY)
    add_title_bar(slide, data['title'], data.get('subtitle', ''))

    tbl_data = data['table']
    if not tbl_data:
        build_content_slide(prs, data, num)
        return

    headers = tbl_data['headers']
    rows = tbl_data['rows']
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)

    table_top = Inches(1.9) if data.get('subtitle') else BODY_TOP
    table_h = SLIDE_H - table_top - Inches(0.8)
    row_h = min(Inches(0.55), table_h / n_rows)

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        MARGIN, table_top, CONTENT_W, row_h * n_rows
    )
    table = tbl_shape.table

    # Remove default table style borders by setting no style
    tbl_el = table._tbl
    tbl_pr = tbl_el.tblPr if tbl_el.tblPr is not None else tbl_el._add_tblPr()
    tbl_pr.set('firstRow', '1')
    tbl_pr.set('bandRow', '1')

    # Style header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = CLR_GOLD
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x15, 0x22, 0x3A)

    # Style data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = CLR_WHITE
            # Alternating row colors
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_DARK_NAVY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CLR_ROW_ALT

            # Color severity cells
            val_lower = val.lower().strip()
            if val_lower == 'critical':
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.color.rgb = CLR_RED
                    paragraph.font.bold = True
            elif val_lower == 'high':
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.color.rgb = CLR_ORANGE
                    paragraph.font.bold = True
            elif val_lower.startswith('go'):
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.color.rgb = CLR_GREEN
            elif val_lower.startswith('conditional'):
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.color.rgb = CLR_ORANGE

    add_slide_number(slide, num)
    add_speaker_notes(slide, data['speaker_notes'])
    add_deep_dive(slide, data.get('deep_dive', ''))


def build_two_column_slide(prs, data, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK_NAVY)
    add_title_bar(slide, data['title'], data.get('subtitle', ''))

    body_top = Inches(1.9) if data.get('subtitle') else BODY_TOP
    col_w = (CONTENT_W - Inches(0.3)) / 2
    body_h = SLIDE_H - body_top - Inches(0.8)

    # Left column
    if data.get('left_column'):
        add_bullets(slide, MARGIN, body_top, col_w, body_h,
                    data['left_column'][:6], font_size=16)

    # Vertical divider line
    mid_x = MARGIN + col_w + Inches(0.1)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        mid_x, body_top, Inches(0.02), body_h
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CLR_GOLD
    line.line.fill.background()

    # Right column
    right_left = mid_x + Inches(0.2)
    if data.get('right_column'):
        add_bullets(slide, right_left, body_top, col_w, body_h,
                    data['right_column'][:6], font_size=16)

    add_slide_number(slide, num)
    add_speaker_notes(slide, data['speaker_notes'])
    add_deep_dive(slide, data.get('deep_dive', ''))


# ── Visual Slide Builders ─────────────────────────────────────────────────────

def _add_shape_box(slide, left, top, width, height, text, fill_color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, font_size=10):
    """Add a colored shape with centered text."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = CLR_WHITE
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = CLR_WHITE
    p.alignment = PP_ALIGN.CENTER
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    return shape


def _add_arrow_label(slide, left, top, text, font_size=8):
    """Add a small label for an arrow/connector."""
    add_textbox(slide, left, top, Inches(1.2), Inches(0.3),
                text, font_size=font_size, color=CLR_GRAY, align=PP_ALIGN.CENTER)


def _add_arrow_shape(slide, left, top, width):
    """Add a right-pointing arrow shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CLR_GRAY
    shape.line.fill.background()
    return shape


def build_visual_solution_flow(prs, data, num):
    """Slide: Our Solution — 6-box left-to-right flow with feedback loop."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK_NAVY)
    add_title_bar(slide, data['title'], data.get('subtitle', ''))

    # 6 boxes in a row
    boxes = [
        ("Process\nHistorian", CLR_ORANGE, MSO_SHAPE.FLOWCHART_MAGNETIC_DISK),
        ("Reflex\nEdge Agent", CLR_GREEN, MSO_SHAPE.ROUNDED_RECTANGLE),
        ("Excel LP\nModel", CLR_BLUE, MSO_SHAPE.ROUNDED_RECTANGLE),
        ("Claude\nTranslation", CLR_PURPLE, MSO_SHAPE.ROUNDED_RECTANGLE),
        ("Slack /\nTeams", CLR_BLUE, MSO_SHAPE.ROUNDED_RECTANGLE),
        ("Operator\n(Decides)", CLR_GRAY, MSO_SHAPE.OVAL),
    ]
    box_w = Inches(1.6)
    box_h = Inches(1.0)
    start_x = Inches(0.6)
    y = Inches(2.8)
    gap = Inches(0.35)

    shape_centers = []
    for i, (label, color, shape_type) in enumerate(boxes):
        x = start_x + i * (box_w + gap)
        _add_shape_box(slide, x, y, box_w, box_h, label, color, shape_type, font_size=10)
        shape_centers.append((x + box_w, y + box_h / 2))  # right edge center

        # Arrow between boxes
        if i < len(boxes) - 1:
            ax = x + box_w
            _add_arrow_shape(slide, ax, y + box_h / 2 - Inches(0.1), gap)

    # Feedback arrow (curved text label below)
    feedback_x = start_x + 5 * (box_w + gap) + Inches(0.2)
    feedback_y = y + box_h + Inches(0.3)
    add_textbox(
        slide, start_x + Inches(1), feedback_y, Inches(9), Inches(0.4),
        "Feedback Loop: Constraint Registry — operator says 'no', system remembers and stops nagging",
        font_size=11, color=CLR_GOLD, align=PP_ALIGN.CENTER, italic=True
    )

    # Trigger icons below
    trigger_y = feedback_y + Inches(0.5)
    add_textbox(slide, Inches(2), trigger_y, Inches(4), Inches(0.3),
                "Process Trigger: equipment drifts outside window",
                font_size=10, color=CLR_GREEN)
    add_textbox(slide, Inches(7), trigger_y, Inches(4), Inches(0.3),
                "Price Trigger: crack spread moves materially",
                font_size=10, color=CLR_GOLD)

    # Bullets below title
    if data['bullets']:
        bullet_y = Inches(5.3)
        add_bullets(slide, MARGIN, bullet_y, CONTENT_W, Inches(1.5),
                    data['bullets'][:4], font_size=14, color=CLR_GRAY)

    add_slide_number(slide, num)
    add_speaker_notes(slide, data['speaker_notes'])
    add_deep_dive(slide, data.get('deep_dive', ''))


def build_visual_system_context(prs, data, num):
    """C4 Level 1 System Context diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK_NAVY)
    add_title_bar(slide, data['title'], data.get('subtitle', ''))

    box_w = Inches(1.5)
    box_h = Inches(0.8)

    # Left column: Data Sources
    ds_x = Inches(0.5)
    ds_labels = [
        ("PI Historian\nProcess Data", CLR_ORANGE, MSO_SHAPE.FLOWCHART_MAGNETIC_DISK),
        ("Market Data\nEIA / OPIS", CLR_BLUE, MSO_SHAPE.DIAMOND),
        ("Excel LP\nModels", CLR_BLUE, MSO_SHAPE.DIAMOND),
    ]
    # Dashed box label
    add_textbox(slide, ds_x, Inches(1.8), Inches(1.8), Inches(0.3),
                "Data Sources", font_size=10, color=CLR_GRAY, align=PP_ALIGN.CENTER)

    ds_positions = []
    for i, (label, color, shape_type) in enumerate(ds_labels):
        y = Inches(2.2) + i * Inches(1.1)
        _add_shape_box(slide, ds_x, y, box_w, box_h, label, color, shape_type, font_size=9)
        ds_positions.append((ds_x + box_w, y + box_h / 2))

    # Center: Reflex Platform
    cx = Inches(4.2)
    cy = Inches(2.8)
    cw = Inches(2.2)
    ch = Inches(1.4)
    reflex = _add_shape_box(slide, cx, cy, cw, ch,
                            "Reflex\nPlatform", CLR_GREEN, font_size=14)
    reflex.line.width = Pt(3)

    # Arrows from data sources to Reflex
    for (sx, sy) in ds_positions:
        _add_arrow_shape(slide, sx + Inches(0.1), sy - Inches(0.1), cx - sx - Inches(0.2))

    # Right column: Delivery channels
    dc_x = Inches(8.0)
    dc_labels = [
        ("Microsoft\nTeams", CLR_BLUE),
        ("Slack", CLR_BLUE),
        ("Email", CLR_BLUE),
    ]
    add_textbox(slide, dc_x, Inches(1.8), Inches(1.8), Inches(0.3),
                "Delivery Channels", font_size=10, color=CLR_GRAY, align=PP_ALIGN.CENTER)

    dc_positions = []
    for i, (label, color) in enumerate(dc_labels):
        y = Inches(2.2) + i * Inches(1.1)
        _add_shape_box(slide, dc_x, y, box_w, box_h, label, color, font_size=9)
        dc_positions.append((dc_x, y + box_h / 2))

    # Arrows from Reflex to delivery
    for (dx, dy) in dc_positions:
        _add_arrow_shape(slide, cx + cw + Inches(0.1), dy - Inches(0.1),
                         dx - cx - cw - Inches(0.2))

    # Claude API below Reflex
    claude_x = Inches(5.0)
    claude_y = Inches(4.8)
    _add_shape_box(slide, claude_x, claude_y, Inches(1.5), Inches(0.6),
                   "Claude API\n(formatting only)", CLR_PURPLE, font_size=9)

    # Far right: Actors
    act_x = Inches(10.5)
    actors = ["Shift\nSupervisors", "LP\nPlanners", "Plant\nManagement"]
    for i, label in enumerate(actors):
        y = Inches(2.2) + i * Inches(1.1)
        _add_shape_box(slide, act_x, y, Inches(1.3), box_h, label, CLR_GRAY,
                       MSO_SHAPE.OVAL, font_size=9)

    # Arrows from delivery to actors
    for i in range(3):
        dy = Inches(2.2) + i * Inches(1.1) + box_h / 2
        _add_arrow_shape(slide, dc_x + box_w + Inches(0.1), dy - Inches(0.1),
                         act_x - dc_x - box_w - Inches(0.2))

    if data['bullets']:
        add_bullets(slide, MARGIN, Inches(5.8), CONTENT_W, Inches(1.2),
                    data['bullets'][:3], font_size=12, color=CLR_GRAY)

    add_slide_number(slide, num)
    add_speaker_notes(slide, data['speaker_notes'])
    add_deep_dive(slide, data.get('deep_dive', ''))


def build_visual_system_arch(prs, data, num):
    """System Architecture Overview — simplified."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK_NAVY)
    add_title_bar(slide, data['title'], data.get('subtitle', ''))

    # Three zone boxes
    zone_y = Inches(2.0)
    zone_h = Inches(3.8)

    # Zone 1: Refinery DMZ
    z1_x = Inches(0.4)
    z1_w = Inches(3.2)
    z1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, z1_x, zone_y, z1_w, zone_h)
    z1.fill.solid()
    z1.fill.fore_color.rgb = RGBColor(0x1F, 0x30, 0x50)
    z1.line.color.rgb = CLR_ORANGE
    z1.line.width = Pt(2)
    # Use dashed line style via XML
    ln = z1.line._ln
    prstDash = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
    ln.append(prstDash)
    add_textbox(slide, z1_x + Inches(0.1), zone_y + Inches(0.05), z1_w - Inches(0.2), Inches(0.3),
                "Refinery DMZ (Purdue 3.5)", font_size=10, color=CLR_ORANGE, bold=True)

    _add_shape_box(slide, z1_x + Inches(0.2), zone_y + Inches(0.5), Inches(2.6), Inches(0.6),
                   "PI Web API Reader", CLR_ORANGE, font_size=9)
    _add_shape_box(slide, z1_x + Inches(0.2), zone_y + Inches(1.3), Inches(2.6), Inches(0.6),
                   "Data Quality Gateway", CLR_GREEN, font_size=9)
    _add_shape_box(slide, z1_x + Inches(0.2), zone_y + Inches(2.1), Inches(2.6), Inches(0.6),
                   "SQLite Buffer", CLR_BLUE, font_size=9)
    add_textbox(slide, z1_x + Inches(0.3), zone_y + Inches(2.9), Inches(2.4), Inches(0.3),
                "HTTPS outbound only", font_size=9, color=CLR_GREEN, italic=True, align=PP_ALIGN.CENTER)

    # Firewall indicator
    fw_x = z1_x + z1_w + Inches(0.1)
    _add_shape_box(slide, fw_x, zone_y + Inches(1.2), Inches(0.5), Inches(1.2),
                   "FW", CLR_RED, MSO_SHAPE.HEXAGON, font_size=8)

    # Zone 2: Reflex Cloud
    z2_x = fw_x + Inches(0.7)
    z2_w = Inches(5.0)
    z2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, z2_x, zone_y, z2_w, zone_h)
    z2.fill.solid()
    z2.fill.fore_color.rgb = RGBColor(0x1F, 0x30, 0x50)
    z2.line.color.rgb = CLR_BLUE
    z2.line.width = Pt(2)
    add_textbox(slide, z2_x + Inches(0.1), zone_y + Inches(0.05), z2_w - Inches(0.2), Inches(0.3),
                "Reflex Cloud (Azure)", font_size=10, color=CLR_BLUE, bold=True)

    modules = [
        ("Data Ingestion", Inches(0.2), Inches(0.5)),
        ("Trigger Engine", Inches(0.2), Inches(1.2)),
        ("LP Orchestrator", Inches(0.2), Inches(1.9)),
        ("AI Translation", Inches(2.6), Inches(0.5)),
        ("Messaging Svc", Inches(2.6), Inches(1.2)),
        ("Constraint Reg.", Inches(2.6), Inches(1.9)),
    ]
    for label, dx, dy in modules:
        _add_shape_box(slide, z2_x + dx, zone_y + dy, Inches(2.2), Inches(0.55),
                       label, CLR_GREEN, font_size=9)

    # DB at bottom of cloud zone
    _add_shape_box(slide, z2_x + Inches(1.2), zone_y + Inches(2.7), Inches(2.4), Inches(0.7),
                   "PostgreSQL +\nTimescaleDB", CLR_ORANGE,
                   MSO_SHAPE.FLOWCHART_MAGNETIC_DISK, font_size=9)

    # Windows LP Worker outside
    wlp_x = z2_x + z2_w + Inches(0.2)
    _add_shape_box(slide, wlp_x, zone_y + Inches(1.5), Inches(1.5), Inches(0.8),
                   "Windows LP\nWorker (Excel)", CLR_RED, font_size=9)

    # Zone 3: Output channels
    out_y = zone_y + Inches(0.3)
    out_x = wlp_x
    out_labels = ["Teams", "Slack", "Email"]
    for i, label in enumerate(out_labels):
        _add_shape_box(slide, wlp_x + Inches(0.1), zone_y + Inches(2.8) + i * Inches(0.35),
                       Inches(1.0), Inches(0.3), label, CLR_BLUE, font_size=8)

    # Dashboard
    add_textbox(slide, z2_x + Inches(0.5), zone_y + Inches(3.5), Inches(3), Inches(0.25),
                "Next.js Dashboard: Operations | Analytics | Admin",
                font_size=9, color=CLR_GRAY, italic=True)

    add_slide_number(slide, num)
    add_speaker_notes(slide, data['speaker_notes'])
    add_deep_dive(slide, data.get('deep_dive', ''))


def build_visual_data_flow(prs, data, num):
    """Data Flow Happy Path — simplified 6-stage pipeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK_NAVY)
    add_title_bar(slide, data['title'], data.get('subtitle', ''))

    # 6 stages as a horizontal pipeline
    stages = [
        ("PI Historian\n+ Edge Agent", CLR_ORANGE, "sensor data\n+ quality check"),
        ("Data Quality\nGateway", CLR_GREEN, "staleness, range,\nmode detection"),
        ("Trigger Engine\n~200 rules", CLR_GREEN, "mode-gated,\ndebounced"),
        ("LP Orchestrator\n+ Excel Worker", CLR_BLUE, "5-min timeout\nwatchdog"),
        ("AI Translation\n+ Claude API", CLR_PURPLE, "numbers by code\nwords by Claude"),
        ("Messaging\nService", CLR_GREEN, "Teams / Slack\n/ Email"),
    ]

    box_w = Inches(1.7)
    box_h = Inches(1.0)
    start_x = Inches(0.4)
    y = Inches(2.5)
    gap = Inches(0.3)

    for i, (label, color, desc) in enumerate(stages):
        x = start_x + i * (box_w + gap)
        _add_shape_box(slide, x, y, box_w, box_h, label, color, font_size=9)
        # Description below
        add_textbox(slide, x, y + box_h + Inches(0.05), box_w, Inches(0.5),
                    desc, font_size=8, color=CLR_GRAY, align=PP_ALIGN.CENTER)
        # Arrow
        if i < len(stages) - 1:
            _add_arrow_shape(slide, x + box_w, y + box_h / 2 - Inches(0.1), gap)

    # Supervisor at end
    sup_x = start_x + 6 * (box_w + gap) - Inches(0.3)
    _add_shape_box(slide, sup_x, y + Inches(0.1), Inches(1.1), Inches(0.8),
                   "Shift\nSupervisor", CLR_GRAY, MSO_SHAPE.OVAL, font_size=9)

    # Key callout
    add_textbox(slide, Inches(1), Inches(4.5), Inches(10), Inches(0.4),
                "End-to-end: ~2-5 minutes from sensor shift to supervisor notification. "
                "Numbers NEVER generated by LLM.",
                font_size=12, color=CLR_GOLD, align=PP_ALIGN.CENTER, bold=True)

    if data['bullets']:
        add_bullets(slide, MARGIN, Inches(5.1), CONTENT_W, Inches(1.5),
                    data['bullets'][:3], font_size=12, color=CLR_GRAY)

    add_slide_number(slide, num)
    add_speaker_notes(slide, data['speaker_notes'])
    add_deep_dive(slide, data.get('deep_dive', ''))


def build_visual_feedback_path(prs, data, num):
    """Feedback Path — supervisor constraint flow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK_NAVY)
    add_title_bar(slide, data['title'], data.get('subtitle', ''))

    y1 = Inches(2.2)
    box_w = Inches(1.8)
    box_h = Inches(0.9)

    # Row 1: Supervisor -> 5-Tap Wizard -> Constraint Registry
    _add_shape_box(slide, Inches(0.5), y1, Inches(1.3), box_h,
                   "Shift\nSupervisor", CLR_GRAY, MSO_SHAPE.OVAL, font_size=9)
    _add_arrow_shape(slide, Inches(1.85), y1 + Inches(0.35), Inches(0.3))

    _add_shape_box(slide, Inches(2.2), y1, Inches(2.5), box_h,
                   "5-Tap Wizard\n1. Unit  2. Type  3. Constraint\n4. Severity  5. Confirm",
                   CLR_GREEN, font_size=9)
    _add_arrow_shape(slide, Inches(4.75), y1 + Inches(0.35), Inches(0.3))

    _add_shape_box(slide, Inches(5.1), y1, box_w, box_h,
                   "Feedback\nProcessor", CLR_GREEN, font_size=9)
    _add_arrow_shape(slide, Inches(6.95), y1 + Inches(0.35), Inches(0.3))

    _add_shape_box(slide, Inches(7.3), y1, Inches(2.2), box_h,
                   "Constraint Registry\n(by equipment, 72h expiry)",
                   CLR_ORANGE, MSO_SHAPE.FLOWCHART_MAGNETIC_DISK, font_size=9)

    # Row 2: Re-solve path
    y2 = Inches(3.8)
    add_textbox(slide, Inches(7.5), y1 + box_h + Inches(0.05), Inches(1.8), Inches(0.3),
                "Re-solve triggered", font_size=9, color=CLR_GOLD,
                align=PP_ALIGN.CENTER, italic=True)

    _add_shape_box(slide, Inches(2.2), y2, box_w, box_h,
                   "LP Orchestrator\n(with new bounds)", CLR_BLUE, font_size=9)
    _add_arrow_shape(slide, Inches(4.05), y2 + Inches(0.35), Inches(0.3))

    _add_shape_box(slide, Inches(4.4), y2, box_w, box_h,
                   "AI Translation\nService", CLR_PURPLE, font_size=9)
    _add_arrow_shape(slide, Inches(6.25), y2 + Inches(0.35), Inches(0.3))

    _add_shape_box(slide, Inches(6.6), y2, box_w, box_h,
                   "Messaging\nService", CLR_GREEN, font_size=9)
    _add_arrow_shape(slide, Inches(8.45), y2 + Inches(0.35), Inches(0.3))

    _add_shape_box(slide, Inches(8.8), y2, Inches(1.5), box_h,
                   "Supervisor\n(revised rec)", CLR_GRAY, MSO_SHAPE.OVAL, font_size=9)

    # Callout box
    callout_y = Inches(5.1)
    callout = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), callout_y, Inches(8), Inches(0.6)
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = RGBColor(0x30, 0x15, 0x15)
    callout.line.color.rgb = CLR_RED
    callout.line.width = Pt(2)
    ln = callout.line._ln
    prstDash = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
    ln.append(prstDash)
    tf = callout.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "System will NOT re-recommend Unit 2 until constraint expires or is manually cleared"
    p.font.size = Pt(12)
    p.font.color.rgb = CLR_RED
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    if data['bullets']:
        add_bullets(slide, MARGIN, Inches(5.9), CONTENT_W, Inches(1.0),
                    data['bullets'][:3], font_size=12, color=CLR_GRAY)

    add_slide_number(slide, num)
    add_speaker_notes(slide, data['speaker_notes'])
    add_deep_dive(slide, data.get('deep_dive', ''))


def build_visual_network(prs, data, num):
    """Network Architecture — OT/IT boundary with Purdue zones."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK_NAVY)
    add_title_bar(slide, data['title'], data.get('subtitle', ''))

    zone_y = Inches(2.0)
    zone_h = Inches(3.5)

    # Zone 1: OT Network (red)
    z1_x = Inches(0.3)
    z1_w = Inches(3.5)
    z1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, z1_x, zone_y, z1_w, zone_h)
    z1.fill.solid()
    z1.fill.fore_color.rgb = RGBColor(0x2A, 0x15, 0x15)
    z1.line.color.rgb = CLR_RED
    z1.line.width = Pt(2)
    ln = z1.line._ln
    prstDash = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
    ln.append(prstDash)
    add_textbox(slide, z1_x + Inches(0.1), zone_y + Inches(0.05), z1_w, Inches(0.3),
                "Purdue Level 0-2: OT Network", font_size=10, color=CLR_RED, bold=True)

    _add_shape_box(slide, z1_x + Inches(0.2), zone_y + Inches(0.5), Inches(1.4), Inches(0.7),
                   "DCS / PLCs\nControl", CLR_RED, font_size=9)
    _add_shape_box(slide, z1_x + Inches(1.8), zone_y + Inches(0.5), Inches(1.4), Inches(0.7),
                   "Sensors\n100+ Tags", CLR_RED, font_size=9)
    _add_shape_box(slide, z1_x + Inches(0.5), zone_y + Inches(1.6), Inches(2.2), Inches(0.7),
                   "PI Data Archive\n(Primary)", CLR_ORANGE,
                   MSO_SHAPE.FLOWCHART_MAGNETIC_DISK, font_size=9)
    add_textbox(slide, z1_x + Inches(0.3), zone_y + Inches(2.5), Inches(2.6), Inches(0.3),
                "Reflex NEVER touches this", font_size=10, color=CLR_RED,
                bold=True, align=PP_ALIGN.CENTER)

    # Firewall 1
    fw1_x = z1_x + z1_w + Inches(0.15)
    _add_shape_box(slide, fw1_x, zone_y + Inches(1.2), Inches(0.5), Inches(1.0),
                   "FW", CLR_RED, MSO_SHAPE.HEXAGON, font_size=8)

    # Zone 2: DMZ (orange)
    z2_x = fw1_x + Inches(0.7)
    z2_w = Inches(3.5)
    z2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, z2_x, zone_y, z2_w, zone_h)
    z2.fill.solid()
    z2.fill.fore_color.rgb = RGBColor(0x2A, 0x20, 0x10)
    z2.line.color.rgb = CLR_ORANGE
    z2.line.width = Pt(2)
    ln2 = z2.line._ln
    prstDash2 = ln2.makeelement(qn('a:prstDash'), {'val': 'dash'})
    ln2.append(prstDash2)
    add_textbox(slide, z2_x + Inches(0.1), zone_y + Inches(0.05), z2_w, Inches(0.3),
                "Purdue Level 3.5: DMZ", font_size=10, color=CLR_ORANGE, bold=True)

    _add_shape_box(slide, z2_x + Inches(0.2), zone_y + Inches(0.5), Inches(2.8), Inches(0.6),
                   "PI Archive — Read-Only Replica", CLR_ORANGE,
                   MSO_SHAPE.FLOWCHART_MAGNETIC_DISK, font_size=9)
    _add_shape_box(slide, z2_x + Inches(0.2), zone_y + Inches(1.3), Inches(2.8), Inches(0.6),
                   "PI Web API (HTTPS:443)", CLR_BLUE, font_size=9)
    _add_shape_box(slide, z2_x + Inches(0.2), zone_y + Inches(2.1), Inches(2.8), Inches(0.8),
                   "Reflex Edge Agent\nPython Docker ~60MB\n+ SQLite Buffer", CLR_GREEN, font_size=9)

    # Firewall 2
    fw2_x = z2_x + z2_w + Inches(0.15)
    _add_shape_box(slide, fw2_x, zone_y + Inches(1.2), Inches(0.5), Inches(1.0),
                   "FW", CLR_RED, MSO_SHAPE.HEXAGON, font_size=8)

    # Zone 3: Cloud (green)
    z3_x = fw2_x + Inches(0.7)
    z3_w = Inches(3.0)
    z3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, z3_x, zone_y, z3_w, zone_h)
    z3.fill.solid()
    z3.fill.fore_color.rgb = RGBColor(0x10, 0x2A, 0x15)
    z3.line.color.rgb = CLR_GREEN
    z3.line.width = Pt(2)
    add_textbox(slide, z3_x + Inches(0.1), zone_y + Inches(0.05), z3_w, Inches(0.3),
                "Purdue Level 4-5: Cloud", font_size=10, color=CLR_GREEN, bold=True)

    _add_shape_box(slide, z3_x + Inches(0.2), zone_y + Inches(0.5), Inches(2.4), Inches(0.6),
                   "Azure IoT Hub", CLR_BLUE, font_size=9)
    _add_shape_box(slide, z3_x + Inches(0.2), zone_y + Inches(1.3), Inches(2.4), Inches(0.6),
                   "Reflex Cloud\nFastAPI Monolith", CLR_GREEN, font_size=9)
    _add_shape_box(slide, z3_x + Inches(0.2), zone_y + Inches(2.1), Inches(2.4), Inches(0.6),
                   "PostgreSQL +\nTimescaleDB", CLR_ORANGE,
                   MSO_SHAPE.FLOWCHART_MAGNETIC_DISK, font_size=9)

    # HTTPS outbound label
    add_textbox(slide, z2_x + z2_w - Inches(0.5), zone_y + Inches(3.1), Inches(2.5), Inches(0.3),
                "HTTPS outbound ONLY", font_size=10, color=CLR_GREEN,
                bold=True, align=PP_ALIGN.CENTER)

    # Bottom callout
    add_textbox(slide, Inches(2), Inches(5.8), Inches(8), Inches(0.3),
                "No inbound firewall rules required. Compatible with hardware data diodes.",
                font_size=11, color=CLR_RED, italic=True, align=PP_ALIGN.CENTER)

    if data['bullets']:
        add_bullets(slide, MARGIN, Inches(6.2), CONTENT_W, Inches(0.8),
                    data['bullets'][:3], font_size=11, color=CLR_GRAY)

    add_slide_number(slide, num)
    add_speaker_notes(slide, data['speaker_notes'])
    add_deep_dive(slide, data.get('deep_dive', ''))


def build_visual_timeline(prs, data, num):
    """Implementation Roadmap — 4-phase horizontal timeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK_NAVY)
    add_title_bar(slide, data['title'], data.get('subtitle', ''))

    # Timeline axis
    axis_y = Inches(3.0)
    axis_left = Inches(0.5)
    axis_right = Inches(12.5)
    axis_w = axis_right - axis_left

    # Thin axis line
    axis_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        axis_left, axis_y, axis_w, Inches(0.03)
    )
    axis_line.fill.solid()
    axis_line.fill.fore_color.rgb = CLR_GRAY
    axis_line.line.fill.background()

    # Phase blocks
    phases = [
        ("Phase 0\nValidation", "Wks 1-8", CLR_ORANGE, Inches(0.5), Inches(1.8),
         "LP Survey | COM Test\nDomain Expert", "Go/No-Go"),
        ("Phase 1\nMVP", "Mo 3-8", CLR_GREEN, Inches(2.5), Inches(3.2),
         "Edge Agent | Triggers | LP\nAI + Messaging | Dashboard", "1 Site\nShadow Mode"),
        ("Phase 2\nBeta", "Mo 9-14", CLR_BLUE, Inches(5.9), Inches(2.5),
         "Active Mode | Analytics\n3-5 Sites | Case Study", "ROI Proven\n$500K+"),
        ("Phase 3\nGA", "Mo 15-24", CLR_PURPLE, Inches(8.6), Inches(3.5),
         "SOC 2 | Multi-Site\nLP Migration | 8-12 Sites", "$600K-$1.5M\nARR"),
    ]

    bar_h = Inches(0.6)
    for label, time_label, color, x, w, desc, milestone in phases:
        # Phase bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, axis_y - bar_h / 2, w, bar_h
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        tf = bar.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(10)
        p.font.color.rgb = CLR_WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Time label above
        add_textbox(slide, x, axis_y - bar_h / 2 - Inches(0.3), w, Inches(0.25),
                    time_label, font_size=9, color=CLR_GRAY, align=PP_ALIGN.CENTER)

        # Description below
        add_textbox(slide, x, axis_y + bar_h / 2 + Inches(0.1), w, Inches(0.6),
                    desc, font_size=9, color=CLR_GRAY, align=PP_ALIGN.CENTER)

        # Milestone diamond at end
        diamond_x = x + w - Inches(0.15)
        diamond = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND,
            diamond_x, axis_y + bar_h / 2 + Inches(0.75), Inches(0.3), Inches(0.3)
        )
        diamond.fill.solid()
        diamond.fill.fore_color.rgb = color
        diamond.line.fill.background()
        add_textbox(slide, diamond_x - Inches(0.4), axis_y + bar_h / 2 + Inches(1.05),
                    Inches(1.2), Inches(0.4),
                    milestone, font_size=8, color=color, bold=True, align=PP_ALIGN.CENTER)

    # "YOU ARE HERE" arrow
    add_textbox(slide, Inches(0.3), Inches(1.8), Inches(2), Inches(0.3),
                "YOU ARE HERE", font_size=12, color=CLR_RED, bold=True,
                align=PP_ALIGN.CENTER)
    _add_arrow_shape(slide, Inches(0.7), Inches(2.1), Inches(0.5))

    # Team size bar at bottom
    add_textbox(slide, Inches(0.5), Inches(5.8), Inches(3), Inches(0.25),
                "2-3 engineers", font_size=9, color=CLR_GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(5.9), Inches(5.8), Inches(2.5), Inches(0.25),
                "3-4 engineers", font_size=9, color=CLR_GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.6), Inches(5.8), Inches(3.5), Inches(0.25),
                "4-6 engineers + 1 CS", font_size=9, color=CLR_GRAY, align=PP_ALIGN.CENTER)

    add_slide_number(slide, num)
    add_speaker_notes(slide, data['speaker_notes'])
    add_deep_dive(slide, data.get('deep_dive', ''))


def build_visual_impl_timeline(prs, data, num):
    """Implementation Timeline (Section 7) — reuse timeline builder."""
    build_visual_timeline(prs, data, num)


# ── Visual Slide Router ──────────────────────────────────────────────────────
# Map slide titles (substring) to custom visual builders
VISUAL_BUILDERS = {
    "Our Solution": build_visual_solution_flow,
    "System Context": build_visual_system_context,
    "System Architecture Overview": build_visual_system_arch,
    "Data Flow — Happy Path": build_visual_data_flow,
    "Data Flow — Feedback Path": build_visual_feedback_path,
    "Network Architecture": build_visual_network,
    "Implementation Roadmap": build_visual_timeline,
    "Implementation Timeline": build_visual_impl_timeline,
}


def build_visual_slide(prs, data, num):
    """Route to a custom visual builder or fall back to bullets + reference."""
    for key, builder in VISUAL_BUILDERS.items():
        if key in data.get('title', ''):
            builder(prs, data, num)
            return

    # Fallback: bullets + diagram reference
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK_NAVY)
    add_title_bar(slide, data['title'], data.get('subtitle', ''))

    if data['bullets']:
        add_bullets(slide, MARGIN, BODY_TOP, CONTENT_W, Inches(3.5),
                    data['bullets'][:6])

    # Diagram reference
    add_textbox(
        slide, Inches(2), Inches(5.5), Inches(8), Inches(0.4),
        "Render at mermaid.live — see docs/design/architecture-diagrams.md for full detail",
        font_size=14, color=CLR_GOLD, align=PP_ALIGN.CENTER, italic=True
    )

    add_slide_number(slide, num)
    add_speaker_notes(slide, data['speaker_notes'])
    add_deep_dive(slide, data.get('deep_dive', ''))


# ── Main Build ────────────────────────────────────────────────────────────────
SLIDE_FILES = [
    "presentation/section1-opportunity-and-risks.md",
    "presentation/section2-changes-and-spec.md",
    "presentation/section3-architecture-visuals.md",
    "presentation/section4-roadmap-and-next.md",
]

TYPE_BUILDERS = {
    'title': build_title_slide,
    'divider': build_divider_slide,
    'content': build_content_slide,
    'table': build_table_slide,
    'two-column': build_two_column_slide,
    'visual': build_visual_slide,
}


def main():
    base = os.path.dirname(os.path.abspath(__file__))
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    all_slides = []
    for f in SLIDE_FILES:
        path = os.path.join(base, f)
        all_slides.extend(parse_slides(path))

    print(f"Parsed {len(all_slides)} slides from {len(SLIDE_FILES)} files")

    for i, data in enumerate(all_slides):
        slide_num = i + 1
        slide_type = data.get('type', 'content')
        builder = TYPE_BUILDERS.get(slide_type, build_content_slide)
        builder(prs, data, slide_num)

    output = os.path.join(base, "Reflex-Team-Walkthrough.pptx")
    prs.save(output)
    print(f"Saved: {output}")
    print(f"Total slides: {len(prs.slides)}")

    # ── Quality Checks ────────────────────────────────────────────────────
    print("\n=== Quality Checks ===")
    total = len(prs.slides)
    print(f"Slide count: {total} {'PASS' if 40 <= total <= 55 else 'FAIL'}")

    # Check speaker notes
    notes_count = 0
    for slide in prs.slides:
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                notes_count += 1
        except Exception:
            pass
    print(f"Slides with speaker notes: {notes_count}/{total}")

    # Check slide numbers (we add textboxes, so count shapes with single-digit/number text)
    numbered = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt.isdigit() and int(txt) <= total:
                    numbered += 1
                    break
    print(f"Slides with slide number: {numbered}/{total}")

    print("=== Done ===")


if __name__ == "__main__":
    main()
